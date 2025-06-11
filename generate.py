from pptx import Presentation
import os
import pandas as pd
import subprocess

input_file = "input.pptx"
output_dir = "output"
pptx_output_file = os.path.join(output_dir, "updated.pptx")
pdf_output_file = os.path.join(output_dir, "updated.pdf")
xlsx_output_file = os.path.join(output_dir, "data.xlsx")

# 1) ÚPRAVA PREZENTACE
prs = Presentation(input_file)
if prs.slides:
    slide = prs.slides[0]
    for shape in slide.shapes:
        if shape.has_text_frame:
            shape.text = "Toto je upravená prezentace"
            break

os.makedirs(output_dir, exist_ok=True)
prs.save(pptx_output_file)

# 2) EXPORT DO PDF pomocí LibreOffice
try:
    subprocess.run([
        "libreoffice",
        "--headless",
        "--convert-to", "pdf",
        "--outdir", output_dir,
        pptx_output_file
    ], check=True)
except Exception as e:
    print(f"Chyba při konverzi do PDF: {e}")

# 3) VÝSTUP DO XLSX (pomocí pandas)
data = {
    "Název": ["Snímek 1"],
    "Text": [slide.shapes[0].text if slide.shapes and slide.shapes[0].has_text_frame else "N/A"]
}
df = pd.DataFrame(data)
df.to_excel(xlsx_output_file, index=False)
